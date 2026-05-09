"""
Text extraction for all supported document formats.

Each extractor returns a list of (section_text, section_number, language) tuples.
section_number is 1-based (page for PDF/PPTX, batch index for CSV/XLSX/DOCX, etc.).
"""
import os
import re
import csv

import fitz  # PyMuPDF
import arabic_reshaper
from langdetect import detect


# ─── Language & text utilities ───────────────────────────────────────────────


def _contains_arabic(text: str) -> bool:
    return bool(re.compile(r'[؀-ۿ]').search(text))


def _fix_arabic_characters(sentence: str) -> str:
    # PDF extraction leaves Arabic in isolated/visual form; reshape restores it
    return arabic_reshaper.reshape(sentence)


def _process_mixed_language(raw_text: str) -> str:
    sentences = re.split(r'(?<=[.!?؟।])\s+|\n', raw_text)
    processed = []
    for sentence in sentences:
        if _contains_arabic(sentence):
            sentence = _fix_arabic_characters(sentence)
        processed.append(sentence)
    return "\n\n".join(processed)


def _detect_language(text: str) -> str:
    try:
        return detect(text)
    except Exception:
        return "en"


def _clean_text(text: str) -> str:
    text = re.sub(r'\n{3,}', '\n\n', text)
    text = re.sub(r' +', ' ', text)
    return text.strip()


def _make_section(raw_text: str, section_number: int) -> tuple[str, int, str]:
    """Clean, normalise Arabic, detect language, and pack into a section tuple."""
    cleaned = _clean_text(_process_mixed_language(raw_text))
    lang    = _detect_language(cleaned)
    return cleaned, section_number, lang


# ─── PDF ─────────────────────────────────────────────────────────────────────


def extract_from_pdf(file_path: str) -> list[tuple]:
    sections = []
    doc = fitz.open(file_path)
    for page_number, page in enumerate(doc):
        blocks = page.get_text("blocks")
        if not blocks:
            continue
        raw_text = "\n\n".join(
            b[4].strip() for b in blocks
            if b[4].strip() and b[6] == 0  # b[6]==0 means text block, not image
        )
        if raw_text.strip():
            sections.append(_make_section(raw_text, page_number + 1))
    doc.close()
    return sections


# ─── TXT ─────────────────────────────────────────────────────────────────────


def extract_from_txt(file_path: str) -> list[tuple]:
    with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
        raw_text = f.read()
    if not raw_text.strip():
        return []
    return [_make_section(raw_text, 1)]


# ─── CSV ─────────────────────────────────────────────────────────────────────

_CSV_BATCH = 50  # rows per section


def extract_from_csv(file_path: str) -> list[tuple]:
    rows_text: list[str] = []
    with open(file_path, 'r', encoding='utf-8-sig', errors='replace', newline='') as f:
        reader  = csv.reader(f)
        headers = next(reader, None)
        for row in reader:
            if headers:
                pairs = [f"{h}: {v}" for h, v in zip(headers, row) if v.strip()]
            else:
                pairs = [v for v in row if v.strip()]
            if pairs:
                rows_text.append(", ".join(pairs))

    sections = []
    for i in range(0, len(rows_text), _CSV_BATCH):
        batch = "\n".join(rows_text[i: i + _CSV_BATCH])
        sections.append(_make_section(batch, i // _CSV_BATCH + 1))
    return sections


# ─── DOCX ────────────────────────────────────────────────────────────────────

_DOCX_BATCH = 30  # paragraphs per section


def extract_from_docx(file_path: str) -> list[tuple]:
    try:
        from docx import Document
    except ImportError:
        raise ImportError("python-docx is required for DOCX files: pip install python-docx")

    doc        = Document(file_path)
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]

    sections = []
    for i in range(0, len(paragraphs), _DOCX_BATCH):
        batch = "\n\n".join(paragraphs[i: i + _DOCX_BATCH])
        sections.append(_make_section(batch, i // _DOCX_BATCH + 1))
    return sections


# ─── DOC (legacy Word) ───────────────────────────────────────────────────────


def extract_from_doc(file_path: str) -> list[tuple]:
    """Tries docx2txt first; falls back to python-docx (some .doc files are OOXML)."""
    try:
        import docx2txt
        raw_text = docx2txt.process(file_path)
        if raw_text and raw_text.strip():
            return [_make_section(raw_text, 1)]
    except (ImportError, Exception):
        pass
    return extract_from_docx(file_path)


# ─── XLSX ────────────────────────────────────────────────────────────────────

_XLSX_BATCH = 50


def extract_from_xlsx(file_path: str) -> list[tuple]:
    try:
        import openpyxl
    except ImportError:
        raise ImportError("openpyxl is required for XLSX files: pip install openpyxl")

    wb          = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
    sections    = []
    section_idx = 1

    for sheet in wb.worksheets:
        headers: list[str] | None = None
        rows_text: list[str]      = []

        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) if c is not None else '' for c in row]
            if headers is None:
                headers = cells
                continue
            pairs = [f"{h}: {v}" for h, v in zip(headers, cells) if v.strip()]
            if pairs:
                rows_text.append(", ".join(pairs))

        for i in range(0, len(rows_text), _XLSX_BATCH):
            batch = f"[Sheet: {sheet.title}]\n" + "\n".join(rows_text[i: i + _XLSX_BATCH])
            sections.append(_make_section(batch, section_idx))
            section_idx += 1

    wb.close()
    return sections


# ─── XLS (legacy Excel) ──────────────────────────────────────────────────────

_XLS_BATCH = 50


def extract_from_xls(file_path: str) -> list[tuple]:
    try:
        import xlrd
    except ImportError:
        raise ImportError("xlrd is required for XLS files: pip install xlrd")

    wb          = xlrd.open_workbook(file_path)
    sections    = []
    section_idx = 1

    for sheet in wb.sheets():
        if sheet.nrows == 0:
            continue
        headers   = [str(sheet.cell_value(0, c)) for c in range(sheet.ncols)]
        rows_text: list[str] = []

        for r in range(1, sheet.nrows):
            pairs = [
                f"{headers[c]}: {sheet.cell_value(r, c)}"
                for c in range(sheet.ncols)
                if str(sheet.cell_value(r, c)).strip()
            ]
            if pairs:
                rows_text.append(", ".join(pairs))

        for i in range(0, len(rows_text), _XLS_BATCH):
            batch = f"[Sheet: {sheet.name}]\n" + "\n".join(rows_text[i: i + _XLS_BATCH])
            sections.append(_make_section(batch, section_idx))
            section_idx += 1

    return sections


# ─── PPTX ────────────────────────────────────────────────────────────────────


def extract_from_pptx(file_path: str) -> list[tuple]:
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError("python-pptx is required for PPTX files: pip install python-pptx")

    prs      = Presentation(file_path)
    sections = []

    for slide_number, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        texts.append(line)
        if texts:
            sections.append(_make_section("\n".join(texts), slide_number))

    return sections


# ─── PPT (legacy PowerPoint) ─────────────────────────────────────────────────


def extract_from_ppt(file_path: str) -> list[tuple]:
    """python-pptx can open some .ppt files; raises a clear error otherwise."""
    try:
        return extract_from_pptx(file_path)
    except Exception as exc:
        raise ValueError(
            f"Cannot parse legacy .ppt file. Convert to .pptx first. Error: {exc}"
        )


# ─── Dispatcher ──────────────────────────────────────────────────────────────

_EXTRACTORS: dict[str, callable] = {
    '.pdf':  extract_from_pdf,
    '.txt':  extract_from_txt,
    '.csv':  extract_from_csv,
    '.docx': extract_from_docx,
    '.doc':  extract_from_doc,
    '.xlsx': extract_from_xlsx,
    '.xls':  extract_from_xls,
    '.pptx': extract_from_pptx,
    '.ppt':  extract_from_ppt,
}

SUPPORTED_EXTENSIONS: frozenset[str] = frozenset(_EXTRACTORS.keys())


def extract_text(file_path: str) -> list[tuple]:
    """
    Dispatch to the correct extractor based on file extension.
    Returns list of (section_text, section_number, language) tuples.
    """
    ext       = os.path.splitext(file_path)[1].lower()
    extractor = _EXTRACTORS.get(ext)
    if extractor is None:
        raise ValueError(
            f"Unsupported file type '{ext}'. "
            f"Supported: {', '.join(sorted(SUPPORTED_EXTENSIONS))}"
        )
    return extractor(file_path)
